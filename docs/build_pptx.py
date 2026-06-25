"""FreiKI Präsentation – Akquise/Sales-Deck (extern)"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Farben ──────────────────────────────────────────────────────────────────
BLUE      = RGBColor(0x1F, 0x54, 0xC0)
NAVY      = RGBColor(0x14, 0x30, 0x6B)
INK       = RGBColor(0x15, 0x29, 0x4A)
MUTED     = RGBColor(0x5A, 0x6B, 0x82)
CANVAS    = RGBColor(0xF4, 0xF6, 0xFA)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GREEN_BG  = RGBColor(0xEA, 0xF7, 0xEE)
GREEN_TXT = RGBColor(0x1F, 0x7A, 0x45)
RED_BG    = RGBColor(0xFD, 0xF1, 0xF0)
RED_TXT   = RGBColor(0xB4, 0x3C, 0x32)
LINE      = RGBColor(0xE3, 0xE8, 0xF0)

W = Inches(13.33)   # Widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # blank

# ── Hilfsfunktionen ──────────────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(BLANK)

def _apply_fill(s, fill, line):
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(0.75)
    else:
        s.line.fill.background()

def rect(slide, x, y, w, h, fill=None, line=None):
    s = slide.shapes.add_shape(1, x, y, w, h)
    _apply_fill(s, fill, line)
    s.text_frame.text = ""
    return s

def txb(slide, text, x, y, w, h,
        size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
        wrap=True, italic=False):
    """Standalone textbox for titles, labels, eyebrows."""
    tf = slide.shapes.add_textbox(x, y, w, h)
    tf.word_wrap = wrap
    p = tf.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tf

def _para(tf, text, size, bold=False, color=INK, align=PP_ALIGN.LEFT,
          space_before=0, space_after=0, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after  = Pt(space_after)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return p

def eyebrow(slide, text, y=Inches(0.55)):
    txb(slide, text.upper(), Inches(0.6), y, Inches(12), Inches(0.35),
        size=10, bold=True, color=BLUE)

def title(slide, text, y=Inches(0.9), size=32):
    txb(slide, text, Inches(0.6), y, Inches(12.1), Inches(1.4),
        size=size, bold=True, color=INK)

def body(slide, text, y=Inches(1.9), size=16, color=MUTED, w=Inches(12.1)):
    txb(slide, text, Inches(0.6), y, w, Inches(4.5),
        size=size, color=color)

def card(slide, x, y, w, h, icon, heading, text, bg=WHITE, border=LINE,
         hcolor=INK, tcolor=MUTED):
    """Single editable shape: emoji + heading + body in one text frame."""
    s = slide.shapes.add_shape(1, x, y, w, h)
    _apply_fill(s, bg, border)
    tf = s.text_frame
    tf.word_wrap = True
    tf.margin_left   = Inches(0.2)
    tf.margin_right  = Inches(0.2)
    tf.margin_top    = Inches(0.2)
    tf.margin_bottom = Inches(0.15)
    _para(tf, icon,    size=20, space_after=6,  first=True)
    _para(tf, heading, size=13, bold=True, color=hcolor, space_after=5)
    _para(tf, text,    size=11, color=tcolor)
    return s

# ═══════════════════════════════════════════════════════════════════════════
# Folie 1 – Titelfolie
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
# Hintergrundverlauf simulieren – oberer blauer Streifen
rect(sl, 0, 0, W, Inches(4.0), fill=NAVY)
rect(sl, 0, Inches(4.0), W, H - Inches(4.0), fill=WHITE)

txb(sl, "Frei", Inches(0.7), Inches(1.2), Inches(5), Inches(1.1),
    size=52, bold=True, color=WHITE)
txb(sl, "KI", Inches(3.15), Inches(1.2), Inches(3), Inches(1.1),
    size=52, bold=True, color=BLUE)

txb(sl, "Der KI-Assistent, der bei Ihnen bleibt",
    Inches(0.7), Inches(2.3), Inches(10), Inches(0.8),
    size=22, color=RGBColor(0xC3, 0xD1, 0xEE))

# Badges
badges = ["🇩🇪 Server in Deutschland", "🔒 DSGVO-konform", "🧠 Lokales LLM", "⚙️ Open Source"]
for i, b in enumerate(badges):
    bx = Inches(0.7) + i * Inches(2.95)
    br = rect(sl, bx, Inches(3.3), Inches(2.7), Inches(0.5),
              fill=GREEN_BG, line=RGBColor(0xCD, 0xE9, 0xD7))
    txb(sl, b, bx + Inches(0.12), Inches(3.34), Inches(2.5), Inches(0.42),
        size=12, bold=True, color=GREEN_TXT)

txb(sl, "Anfragen: admin@freiki.net",
    Inches(0.7), Inches(4.5), Inches(8), Inches(0.5),
    size=14, color=MUTED)

# ═══════════════════════════════════════════════════════════════════════════
# Folie 2 – Die einfache Erklärung
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, W, H, fill=CANVAS)
eyebrow(sl, "Die einfache Erklärung")
title(sl, "Fremdes Büro oder eigener Mitarbeiter?")

# Linke Karte (schlecht)
card(sl, Inches(0.6), Inches(2.0), Inches(5.8), Inches(4.6),
     "☁️", "API-KI (z. B. ChatGPT)",
     "Das fremde Büro. Sie schicken Ihre Frage per Brief an ein Büro im Ausland. "
     "Das Büro liest alles – und schickt eine Antwort zurück. Ob Kopien gemacht "
     "oder Infos weitergegeben werden, wissen Sie nicht.",
     bg=RED_BG, border=RGBColor(0xF4, 0xD3, 0xCF),
     hcolor=RED_TXT, tcolor=INK)
txb(sl, "⚠ Daten verlassen Ihre Organisation",
    Inches(0.82), Inches(5.8), Inches(5.3), Inches(0.45),
    size=12, bold=True, color=RED_TXT)

# Rechte Karte (gut)
card(sl, Inches(6.93), Inches(2.0), Inches(5.8), Inches(4.6),
     "🏢", "FreiKI",
     "Der eigene Mitarbeiter. Sie flüstern Ihre Frage einem Mitarbeiter zu, "
     "der exklusiv für Sie in einem abgesicherten Büro sitzt. Er antwortet sofort "
     "– und verlässt dieses Büro niemals.",
     bg=GREEN_BG, border=RGBColor(0xCD, 0xE9, 0xD7),
     hcolor=GREEN_TXT, tcolor=INK)
txb(sl, "✓ Daten bleiben bei Ihnen",
    Inches(7.15), Inches(5.8), Inches(5.3), Inches(0.45),
    size=12, bold=True, color=GREEN_TXT)

# ═══════════════════════════════════════════════════════════════════════════
# Folie 3 – Was FreiKI kann (Funktionsübersicht)
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
eyebrow(sl, "Funktionen")
title(sl, "Was FreiKI kann")
txb(sl, "Ein Werkzeugkasten statt eines einzelnen Chatfensters – jeder Baustein lässt sich einzeln freischalten.",
    Inches(0.6), Inches(1.75), Inches(12.1), Inches(0.55), size=14, color=MUTED)

features = [
    ("💬", "Chat", "Freies Gespräch mit dem lokalen LLM"),
    ("📚", "Wissen (RAG)", "Suche in eigenen Wissensbereichen"),
    ("🎙️", "Transkription", "Audio → Text → Audio, inkl. Vorlesen"),
    ("🔎", "Websuche", "Interne Recherche via SearXNG"),
    ("🌍", "100+ Sprachen", "Für international zusammengesetzte Teams"),
    ("📤", "Wissen hinterlegen", "Unternehmenswissen einfach befragen"),
    ("📷", "Foto-OCR", "Text aus Fotos erkennen & nutzen"),
    ("📄", "MultiDoc", "Mehrere Dokumente gleichzeitig vergleichen"),
    ("🗣️", "Leichte Sprache", "Texte einfach und verständlich"),
    ("🖼️", "Piktogramme", "Bildkarten (ARASAAC) herunterladen"),
    ("💬", "Team-Chat", "KI-Assistent direkt im Mattermost"),
    ("📄", "DMS", "KI-gestütztes Dokumentenarchiv"),
]
cols = 4
cw = Inches(3.05)
ch = Inches(2.1)
for i, (icon, h, t) in enumerate(features):
    col = i % cols
    row = i // cols
    cx = Inches(0.6) + col * (cw + Inches(0.15))
    cy = Inches(2.4) + row * (ch + Inches(0.12))
    card(sl, cx, cy, cw, ch, icon, h, t)

# ═══════════════════════════════════════════════════════════════════════════
# Folie 4 – Datenschutz & Werkstatt-Bild
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, W, H, fill=CANVAS)
eyebrow(sl, "Datenschutz konkret")
title(sl, "Wie eine Mietwerkstatt – nur für KI")
txb(sl, "Stellen Sie sich FreiKI wie eine Mietwerkstatt vor: Sie nutzen Chat, Wissenssuche, Übersetzung, "
    "Transkription – so viele Teilwerkstätten, wie Sie brauchen. Wenn Sie gehen, weiß niemand, "
    "was Sie darin gemacht haben.",
    Inches(0.6), Inches(1.75), Inches(12.1), Inches(0.9), size=15, color=MUTED)

# Flow
steps = [("💬", "Frage gestellt"), ("🤖", "FreiKI antwortet"), ("🧹", "Verlauf gelöscht"), ("⬜", "Leerer Speicher")]
sw = Inches(2.4)
sh = Inches(1.6)
sx0 = Inches(0.85)
sy = Inches(3.0)
for i, (ic, lab) in enumerate(steps):
    bx = sx0 + i * (sw + Inches(0.65))
    r = rect(sl, bx, sy, sw, sh, fill=WHITE, line=LINE)
    txb(sl, ic, bx, sy + Inches(0.15), sw, Inches(0.7), size=28, align=PP_ALIGN.CENTER)
    txb(sl, lab, bx, sy + Inches(0.9), sw, Inches(0.55),
        size=13, bold=True, color=INK, align=PP_ALIGN.CENTER)
    if i < 3:
        txb(sl, "→", bx + sw + Inches(0.1), sy + Inches(0.55), Inches(0.5), Inches(0.5),
            size=22, color=MUTED, align=PP_ALIGN.CENTER)

txb(sl, "Anfragen werden serverseitig nicht gespeichert. "
    "Der Chatverlauf bleibt ausschließlich in Ihrem Browser erhalten.",
    Inches(0.6), Inches(5.0), Inches(12.1), Inches(0.6),
    size=13, color=MUTED, align=PP_ALIGN.CENTER)

# 3 Datenschutz-Punkte
dp = [
    ("🇩🇪", "Server in Deutschland", "Kein Rechenzentrum in den USA oder bei Hyperscalern."),
    ("🚫", "Kein Training mit Ihren Daten", "Das Modell lernt nicht aus Ihren Eingaben – keine Rückkopplung an den Hersteller."),
    ("🔐", "Keine Cloud-Abhängigkeit", "Alle Komponenten laufen auf Ihrer Infrastruktur. Kein Single Point of Dependency."),
]
for i, (ic, h, t) in enumerate(dp):
    cx = Inches(0.6) + i * Inches(4.2)
    card(sl, cx, Inches(5.8), Inches(3.9), Inches(1.45), ic, h, t)

# ═══════════════════════════════════════════════════════════════════════════
# Folie 5 – Unternehmenswissen (RAG)
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
eyebrow(sl, "Wissen (RAG)")
title(sl, "Unternehmenswissen bewahren und erschließen")
txb(sl, "Handbücher, Richtlinien, interne Dokumentationen – meist verteilt über viele Ordner und Köpfe. "
    "FreiKI macht dieses Wissen durchsuchbar und für alle nutzbar.",
    Inches(0.6), Inches(1.75), Inches(12.1), Inches(0.65), size=15, color=MUTED)

rag = [
    ("🗄️", "Wissen bewahren",
     "Erfahrung und Fachwissen gehen nicht verloren, wenn Mitarbeitende das Unternehmen verlassen "
     "– Dokumente und Prozesse bleiben dauerhaft auffindbar."),
    ("🔍", "Wissen erschließen",
     "Statt mühsam in Ordnerstrukturen zu suchen, stellt man die Frage direkt – FreiKI durchsucht "
     "alle freigegebenen Dokumente und liefert die passende Antwort mit Quellenangabe."),
    ("🤝", "Onboarding & Sprachvielfalt",
     "Neue und internationale Mitarbeitende finden sich schneller zurecht: Fragen zu internen "
     "Abläufen lassen sich jederzeit stellen – in über 100 Sprachen beantwortet."),
]
cw = Inches(4.0)
for i, (ic, h, t) in enumerate(rag):
    cx = Inches(0.6) + i * (cw + Inches(0.22))
    card(sl, cx, Inches(2.55), cw, Inches(4.6), ic, h, t)

# ═══════════════════════════════════════════════════════════════════════════
# Folie 6 – DMS Paperless
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, W, H, fill=CANVAS)
eyebrow(sl, "KI-unterstütztes Dokumentenmanagement")
title(sl, "Paperless-ngx: das digitale Archiv mit KI-Verstand")
txb(sl, "FreiKI integriert Paperless-ngx als vollständig selbst gehostetes DMS. "
    "Dokumente werden eingescannt oder per Mail angeliefert – die KI übernimmt den Rest.",
    Inches(0.6), Inches(1.75), Inches(12.1), Inches(0.65), size=15, color=MUTED)

dms = [
    ("🏷️", "Automatisches Verschlagworten",
     "Eingehende Dokumente werden gelesen, erkannt und automatisch mit Titeln, Tags und "
     "Dokumenttypen versehen – kein manuelles Einsortieren mehr."),
    ("🔍", "Volltextsuche im Archiv",
     "Alle archivierten Dokumente sind sofort durchsuchbar – per Freitext oder direkt "
     "aus dem FreiKI-Chat heraus über die integrierten Wissensbereiche."),
    ("🔒", "Dokumente verlassen Ihr Haus nicht",
     "Weder für die Texterkennung noch für die KI-Analyse wird ein externer Dienst genutzt. "
     "Alles läuft auf Ihrer Infrastruktur – DSGVO-konform und ohne Cloudabhängigkeit."),
]
cw = Inches(4.0)
for i, (ic, h, t) in enumerate(dms):
    cx = Inches(0.6) + i * (cw + Inches(0.22))
    card(sl, cx, Inches(2.55), cw, Inches(4.6), ic, h, t)

# ═══════════════════════════════════════════════════════════════════════════
# Folie 7 – Team-Chat Mattermost
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
eyebrow(sl, "Neu: Team-Chat")
title(sl, "Mattermost – Team-Kommunikation und KI in einem")
txb(sl, "FreiKI bringt einen vollständig selbst gehosteten Team-Chat mit: Mattermost (Open Source). "
    "Teamkommunikation und KI-Unterstützung laufen auf derselben Infrastruktur – "
    "datenschutzkonform und ohne Cloud-Abhängigkeit.",
    Inches(0.6), Inches(1.75), Inches(12.1), Inches(0.75), size=15, color=MUTED)

mm = [
    ("🔒", "Kein Slack, kein Teams – Ihre Daten bleiben bei Ihnen",
     "Bei Slack, Microsoft Teams oder Google Chat liegen Ihre Gespräche auf fremden Servern "
     "– oft in den USA. Mattermost läuft ausschließlich auf Ihrer eigenen Infrastruktur."),
    ("🇩🇪", "DSGVO-konform ohne Kompromisse",
     "Keine Drittanbieter, keine Datenweitergabe, kein AVV mit einem US-Konzern nötig. "
     "Mattermost läuft auf demselben deutschen Server wie FreiKI."),
    ("🤖", "FreiKI-Bot direkt im Chat",
     "Die KI ist mitten im Team-Gespräch erreichbar – Antworten inklusive Quellenangaben "
     "aus den internen Wissensbereichen, ohne dass Ihre Fragen das Haus verlassen."),
]
cw = Inches(4.0)
for i, (ic, h, t) in enumerate(mm):
    cx = Inches(0.6) + i * (cw + Inches(0.22))
    card(sl, cx, Inches(2.65), cw, Inches(4.5), ic, h, t)

# ═══════════════════════════════════════════════════════════════════════════
# Folie 8 – Agenten / Automatisierung
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, W, H, fill=CANVAS)
eyebrow(sl, "Automatisierung")
title(sl, "Mehr als ein Chat: Agenten im Hintergrund")
txb(sl, "Über die eingebaute Automatisierungs-Engine n8n lassen sich eigene KI-Agenten bauen, "
    "die selbstständig arbeiten – ganz ohne Klick in der Oberfläche.",
    Inches(0.6), Inches(1.75), Inches(12.1), Inches(0.65), size=15, color=MUTED)

ag = [
    ("📬", "Abwesenheitsassistent",
     "Liest eingehende Mails, ignoriert Newsletter, leitet wichtige Anfragen weiter und "
     "beantwortet den Rest automatisch mit einer passenden Abwesenheitsinfo."),
    ("📡", "Selbstständige Recherche & Meldungen",
     "Agenten durchsuchen eigenständig das Internet – z. B. nach Wetterwarnungen oder "
     "NINA-Katastrophenschutzmeldungen – und stellen relevante Hinweise per E-Mail oder "
     "Team-Chat zu."),
    ("🚨", "Eskalations-Workflow",
     "Die KI bewertet die Dringlichkeit einer Anfrage und löst bei Bedarf sofort eine "
     "Benachrichtigung aus – statt auf die normale Bearbeitung zu warten."),
]
cw = Inches(4.0)
for i, (ic, h, t) in enumerate(ag):
    cx = Inches(0.6) + i * (cw + Inches(0.22))
    card(sl, cx, Inches(2.55), cw, Inches(4.6), ic, h, t)

# ═══════════════════════════════════════════════════════════════════════════
# Folie 9 – Kosten & Lizenz
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
eyebrow(sl, "Lizenz & Betrieb")
title(sl, "Open Source – Sie selbst sind der Betreiber")
txb(sl, "Alle Bausteine von FreiKI sind frei verfügbar. Sie sind nicht an einen einzelnen Anbieter "
    "gebunden. Die einzigen laufenden Kosten sind Server und GPU.",
    Inches(0.6), Inches(1.75), Inches(12.1), Inches(0.65), size=15, color=MUTED)

# Stats
stats = [("0 €", "Lizenzkosten"), ("100 %", "Daten in eigener Infrastruktur"), ("frei", "wählbares LLM")]
for i, (num, lab) in enumerate(stats):
    sx = Inches(0.8) + i * Inches(4.1)
    txb(sl, num, sx, Inches(2.55), Inches(3.8), Inches(0.8),
        size=36, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    txb(sl, lab, sx, Inches(3.35), Inches(3.8), Inches(0.45),
        size=13, color=MUTED, align=PP_ALIGN.CENTER)

# Preis-Karten
for i, (h, amount, items) in enumerate([
    ("Cloud (EU)", "ca. 1.500 €/Monat", [
        "Schnell startklar, flexibel",
        "Europäische Anbieter (z. B. Hetzner, OVHcloud)",
        "Dedizierter GPU-Server für den Produktivbetrieb",
        "Keine Einstiegskosten",
    ]),
    ("On-Premises", "ab ca. 50.000 € einmalig", [
        "Eigener GPU-Server, einmalige Anschaffung",
        "Danach nur Strom & Netzanbindung",
        "Ab ca. 3 Jahren günstiger als Cloud-Miete",
    ]),
]):
    px = Inches(0.6) + i * Inches(6.3)
    r = rect(sl, px, Inches(4.1), Inches(6.0), Inches(3.0), fill=WHITE, line=LINE)
    txb(sl, h, px + Inches(0.25), Inches(4.2), Inches(5.5), Inches(0.45),
        size=15, bold=True, color=INK)
    txb(sl, amount, px + Inches(0.25), Inches(4.65), Inches(5.5), Inches(0.55),
        size=20, bold=True, color=BLUE)
    for j, item in enumerate(items):
        txb(sl, "✓  " + item,
            px + Inches(0.25), Inches(5.25) + j * Inches(0.42),
            Inches(5.5), Inches(0.4),
            size=12, color=MUTED)

# ═══════════════════════════════════════════════════════════════════════════
# Folie 10 – FAQ (Auswahl)
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, W, H, fill=CANVAS)
eyebrow(sl, "Häufige Fragen")
title(sl, "Was Skeptiker fragen – und unsere Antworten")

faqs = [
    ("Ist das wirklich so gut wie ChatGPT?",
     "Für interne Zwecke: ja. Was Sie gewinnen, ist vollständiger Datenschutz. "
     "Was Sie eventuell aufgeben, sind die allerneuesten Web-Informationen."),
    ("Lernt die KI aus unseren Daten?",
     "Nein. Das Modell ist statisch und trainiert nicht mit Ihren Eingaben. "
     "Es gibt keine Rückkopplung an den Modell-Hersteller."),
    ("Werden meine Anfragen gespeichert?",
     "Nein. Nach jeder Antwort werden Verlauf und Kontext serverseitig vollständig gelöscht. "
     "Der Chatverlauf bleibt nur in Ihrem Browser und nur für die laufende Sitzung."),
    ("Was, wenn der Server ausfällt?",
     "Dann ist FreiKI vorübergehend nicht verfügbar. Professionelle Rechenzentren bieten hohe "
     "Verfügbarkeit – deutlich mehr als ein eigener Server im Haus."),
]
for i, (q, a) in enumerate(faqs):
    fy = Inches(1.85) + i * Inches(1.35)
    r = rect(sl, Inches(0.6), fy, Inches(12.13), Inches(1.2), fill=WHITE, line=LINE)
    txb(sl, "Q  " + q, Inches(0.85), fy + Inches(0.1), Inches(11.6), Inches(0.4),
        size=13, bold=True, color=INK)
    txb(sl, a, Inches(0.85), fy + Inches(0.5), Inches(11.6), Inches(0.62),
        size=12, color=MUTED)

# ═══════════════════════════════════════════════════════════════════════════
# Folie 11 – Abschluss / CTA
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, W, H, fill=NAVY)

txb(sl, "FreiKI", Inches(0.6), Inches(0.5), Inches(12), Inches(0.55),
    size=12, bold=True, color=RGBColor(0x8F, 0xB1, 0xF2), align=PP_ALIGN.CENTER)

txb(sl, "Datenschutz, Unabhängigkeit und volle Kontrolle –\nKI in Ihrer eigenen Organisation.",
    Inches(0.6), Inches(1.3), Inches(12.13), Inches(1.8),
    size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txb(sl, "Bei Interesse Anfrage an admin@freiki.net\nSie erhalten ein Passwort für den Testserver.",
    Inches(0.6), Inches(3.3), Inches(12.13), Inches(0.85),
    size=16, color=RGBColor(0xC3, 0xD1, 0xEE), align=PP_ALIGN.CENTER)

# CTA-Button
btn = rect(sl, Inches(5.0), Inches(4.4), Inches(3.33), Inches(0.65), fill=BLUE)
btn.line.fill.background()
txb(sl, "Zugang anfragen →", Inches(5.0), Inches(4.45), Inches(3.33), Inches(0.55),
    size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

tags = "Deutscher Server  ·  DSGVO/DSG-EKD-konform  ·  Ohne Lizenzrisiken  ·  Open Source"
txb(sl, tags, Inches(0.6), Inches(5.5), Inches(12.13), Inches(0.5),
    size=12, color=RGBColor(0x8F, 0xB1, 0xF2), align=PP_ALIGN.CENTER)

# ── Speichern ───────────────────────────────────────────────────────────────
out = "/Users/frankstefan/ClaudePro/freiki-package/docs/FreiKI-Praesentation.pptx"
prs.save(out)
print("Gespeichert:", out)
print(f"Folien: {len(prs.slides)}")
