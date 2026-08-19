#!/usr/bin/env python3
"""Füllt das Diakonie-Kork-PPTX-Template mit Folien aus JSON (passendes Layout wird
anhand der Platzhalter-Struktur gesucht, nicht anhand des Layout-Namens, siehe
find_content_layout()/find_title_only_layout()). Wird von PptxExportService.js als
Kindprozess aufgerufen: argv[1] = Pfad zur Eingabe-JSON, argv[2] = Zielpfad für die
erzeugte .pptx.

JSON-Form:
{
  "templatePath": "...pptx",
  "slides": [
    {
      "title": "...",
      "body": [{"text": "...", "bullet": true, "number": false, "indent": 0, "bold": false}],
      "image": "/pfad/zu/bild.png"   // optional, verkleinert den Textbereich und
                                       // setzt das Bild rechts daneben
    }
  ]
}
"""
import json
import sys

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt
from PIL import Image

# Das Template deklariert "buNone" auf allen Textebenen (siehe set_bullet() unten),
# Bullets müssen daher pro Absatz manuell gesetzt werden.
IMAGE_MARGIN = 274638  # ~0.3in Abstand zwischen Text- und Bildspalte

# Layouts werden bewusst NICHT über den Namen gesucht: Der Kunde pflegt das Template
# selbst in PowerPoint, und Layout-Namen/-Anzahl haben sich zwischen zwei Versionen
# bereits ohne inhaltliche Absicht geändert (z.B. "Title, Content" -> "Titelfolie").
# Stattdessen wird nach Platzhalter-Struktur gesucht, das bleibt über Rebuilds stabil.
def find_content_layout(prs):
    for layout in prs.slide_layouts:
        phs = {ph.placeholder_format.idx: ph for ph in layout.placeholders}
        title, body = phs.get(0), phs.get(1)
        if title is not None and body is not None and body.width >= prs.slide_width * 0.6:
            return layout
    raise ValueError('Kein Layout mit Titel + durchgehend breitem Textplatzhalter im Template gefunden')


def find_title_only_layout(prs):
    for layout in prs.slide_layouts:
        phs = {ph.placeholder_format.idx: ph for ph in layout.placeholders}
        if 0 in phs and 1 not in phs:
            return layout
    raise ValueError('Kein reines Titel-Layout (Titel ohne Textplatzhalter) im Template gefunden')


def set_bullet(paragraph, level, numbered):
    """pptxgenjs/python-pptx haben keine High-Level-Bullet-API; das Template
    setzt buNone auf Layout-Ebene (freies Design ohne Bullets als Default),
    also injizieren wir buChar/buAutoNum direkt in die pPr-XML des Absatzes."""
    paragraph.level = level
    pPr = paragraph._pPr
    if pPr is None:
        pPr = paragraph._p.get_or_add_pPr()
    for tag in ('a:buNone', 'a:buChar', 'a:buAutoNum'):
        el = pPr.find(qn(tag))
        if el is not None:
            pPr.remove(el)
    marL = 228600 + level * 457200
    pPr.set('marL', str(marL))
    pPr.set('indent', str(-228600))
    if numbered:
        bu = pPr.makeelement(qn('a:buAutoNum'), {'type': 'arabicPeriod'})
    else:
        bu = pPr.makeelement(qn('a:buChar'), {'char': '•'})
    pPr.append(bu)


def fit_image_box(image_path, box):
    left, top, width, height = box
    with Image.open(image_path) as im:
        iw, ih = im.size
    box_ratio = width / height
    img_ratio = iw / ih
    if img_ratio > box_ratio:
        w = width
        h = int(width / img_ratio)
    else:
        h = height
        w = int(height * img_ratio)
    x = left + (width - w) // 2
    y = top + (height - h) // 2
    return x, y, w, h


def add_body_text(placeholder, items):
    tf = placeholder.text_frame
    tf.clear()
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = item.get('text', '')
        run.font.bold = bool(item.get('bold'))
        if item.get('bullet') or item.get('number'):
            set_bullet(p, min(int(item.get('indent') or 0), 4), bool(item.get('number')))
        else:
            # Fließtext ohne Bullet (z.B. Zwischenüberschriften) bekommt buNone.
            p.level = min(int(item.get('indent') or 0), 4)
            pPr = p._p.get_or_add_pPr()
            pPr.append(pPr.makeelement(qn('a:buNone'), {}))


def remove_existing_slides(prs):
    """Die .potx/.pptx-Vorlage selbst enthaelt noch eine (leere) Restfolie aus dem
    Original-Speicherstand - die muss vor dem Befuellen raus, sonst haengt sie als
    zusaetzliche erste Folie im Export."""
    xml_slides = prs.slides._sldIdLst
    for sld_id in list(xml_slides):
        rId = sld_id.get(qn('r:id'))
        prs.part.drop_rel(rId)
        xml_slides.remove(sld_id)


def build(data, out_path):
    prs = Presentation(data['templatePath'])
    remove_existing_slides(prs)
    content_layout = find_content_layout(prs)
    title_only_layout = find_title_only_layout(prs)
    body_ph_def = next(ph for ph in content_layout.placeholders if ph.placeholder_format.idx == 1)
    BODY_FULL = (body_ph_def.left, body_ph_def.top, body_ph_def.width, body_ph_def.height)

    for slide_data in data['slides']:
        image_path = slide_data.get('image')
        has_body = bool(slide_data.get('body'))
        layout = content_layout if has_body else title_only_layout
        slide = prs.slides.add_slide(layout)
        title_ph = slide.placeholders[0]
        title_ph.text_frame.text = slide_data.get('title') or ' '
        # Layout hinterlegt fuer den Titel nur 18pt (siehe Docstring oben) - die
        # eigentliche Titelgroesse ist im echten Template offenbar von Hand pro
        # Folie gesetzt, hier deshalb explizit statt vom Layout geerbt.
        title_run = title_ph.text_frame.paragraphs[0].runs[0]
        title_run.font.size = Pt(40)
        title_run.font.bold = True

        if has_body:
            body_ph = slide.placeholders[1]
            if image_path:
                text_w = (BODY_FULL[2] - IMAGE_MARGIN) // 2
                body_ph.left, body_ph.top = Emu(BODY_FULL[0]), Emu(BODY_FULL[1])
                body_ph.width, body_ph.height = Emu(text_w), Emu(BODY_FULL[3])
                add_body_text(body_ph, slide_data['body'])
                image_box = (BODY_FULL[0] + text_w + IMAGE_MARGIN, BODY_FULL[1],
                             BODY_FULL[2] - text_w - IMAGE_MARGIN, BODY_FULL[3])
                x, y, w, h = fit_image_box(image_path, image_box)
                slide.shapes.add_picture(image_path, Emu(x), Emu(y), width=Emu(w), height=Emu(h))
            else:
                add_body_text(body_ph, slide_data['body'])
        elif image_path:
            x, y, w, h = fit_image_box(image_path, BODY_FULL)
            slide.shapes.add_picture(image_path, Emu(x), Emu(y), width=Emu(w), height=Emu(h))

    prs.save(out_path)


def main():
    in_path, out_path = sys.argv[1], sys.argv[2]
    with open(in_path, encoding='utf-8') as f:
        data = json.load(f)
    build(data, out_path)


if __name__ == '__main__':
    main()
